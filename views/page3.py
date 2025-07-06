import streamlit as st
import re
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from fpdf import FPDF
from openvino.runtime import Core
import speech_recognition as sr
from groq import Groq
from dotenv import load_dotenv
import os


# Load Groq API key from .env
load_dotenv()
GROQ_API_KEY = os.getenv("GROQ_API")
if not GROQ_API_KEY:
    st.error("GROQ_API not set in .env file")
    st.stop()
groq_client = Groq(api_key=GROQ_API_KEY)

# --- Load OpenVINO model ---
@st.cache_resource
def load_openvino_model():
    model_path = "./openvino_models/miniLM_openvino"
    core = Core()
    try:
        compiled_model = core.compile_model(f"{model_path}/openvino_model.xml", "CPU")

        try:
            from transformers import AutoTokenizer
            tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
        except:
            from sentence_transformers import SentenceTransformer
            model = SentenceTransformer('all-MiniLM-L6-v2')
            tokenizer = model.tokenizer

        return compiled_model, tokenizer
    except Exception as e:
        st.error(f"Error loading model: {e}")
        st.stop()

# --- Document Understanding System ---

# --- AI-powered Document Chat using Groq ---
class DocumentUnderstanding:
    def __init__(self):
        self.full_text = ""
        self.sentences = []

    def process_document(self, pages, compiled_model, tokenizer):
        self.full_text = " ".join(pages)
        # Optionally split into sentences for other features
        self.sentences = re.split(r'(?<=[.!?])\s+', self.full_text)

    def answer_question(self, question, compiled_model, tokenizer):
        """Use Groq LLM to answer question based on document content"""
        if not self.full_text.strip():
            return "No document content available to answer questions."
        try:
            prompt = (
                "You are an expert assistant. Answer the user's question using ONLY the following document content. "
                "If the answer is not present, say so.\n\n"
                f"Document Content:\n{self.full_text[:12000]}\n\nQuestion: {question}\nAnswer:"
            )
            chat_completion = groq_client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama-3.1-8b-instant",
            )
            return chat_completion.choices[0].message.content.strip()
        except Exception as e:
            st.error(f"Groq API error: {e}")
            return "Could not generate an answer. Please try a different question."

# --- Text cleaning ---
IGNORE_PHRASES = ['university', 'course', 'module', 'syllabus', 'lecture',
                 'notes', 'professor', 'department', 'b.tech', 'cse', 'git', 'gu']

def clean_text(text):
    text = re.sub(r'[\x00-\x1F]+', ' ', text)
    text = re.sub(r'\b(' + '|'.join(IGNORE_PHRASES) + r')\b', '', text, flags=re.IGNORECASE)
    text = re.sub(r'[^a-zA-Z0-9.,:\-()\n ]+', '', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

# --- File extraction ---
def extract_pdf(file):
    return [page.extract_text() or '' for page in PdfReader(file).pages]

def extract_docx(file):
    return [p.text for p in docx.Document(file).paragraphs if p.text.strip()]

def extract_pptx(file):
    return [shape.text.strip() for slide in Presentation(file).slides
            for shape in slide.shapes
            if hasattr(shape, 'text') and shape.text.strip()]

# --- Notes Generation ---
def generate_quality_notes(text):
    text = clean_text(text)
    sentences = re.split(r'(?<=[.!?])\s+', text)
    notes = []

    for s in sentences:
        s = s.strip()
        if len(s) < 10 or any(phrase in s.lower() for phrase in IGNORE_PHRASES):
            continue

        s = s[0].upper() + s[1:]  # Capitalize first letter
        if not s.endswith(('.','!','?')):
            s += '.'

        notes.append(f"• {s}")

    return notes[:200]

# --- Headings Extraction ---
def extract_headings_and_content(pages):
    content_map = {}
    current_heading = None

    for page in pages:
        lines = page.split('\n') if isinstance(page, str) else [page]
        for line in lines:
            line = line.strip()
            if not line:
                continue

            if (len(line) < 80 and
                (line.endswith(':') or line.isupper() or line.istitle() or
                 line.lower().startswith(('chapter','section','part','unit')))):
                current_heading = line.strip(':')
                if current_heading not in content_map:
                    content_map[current_heading] = ''
            elif current_heading:
                content_map[current_heading] += ' ' + line
            else:
                current_heading = "Key Concepts"
                content_map[current_heading] = ' ' + line

    return content_map

# --- PDF Generation ---
def generate_pdf(notes):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.set_auto_page_break(True, margin=15)

    for n in notes:
        try:
            clean_n = re.sub(r'(\S{30})(?=\S)', r'\1 ', n)
            clean_n = clean_n.encode('latin-1', 'replace').decode('latin-1')
            pdf.multi_cell(180, 10, clean_n)
            pdf.ln(5)
        except Exception as e:
            st.error(f"Error processing note: {e}")
            continue

    return bytes(pdf.output(dest='S').encode('latin1'))

def listen():
    r = sr.Recognizer()
    try:
        with sr.Microphone() as source:
            st.toast("Listening... Speak now (5 second timeout)")
            r.adjust_for_ambient_noise(source)
            audio = r.listen(source, timeout=5, phrase_time_limit=5)
        text = r.recognize_google(audio)
        return text
    except Exception as e:
        st.toast(f"Couldn't understand speech: {e}")
        return None

# --- Streamlit UI ---
st.title("🤝 AI Buddy")

compiled_model, tokenizer = load_openvino_model()

if "library" not in st.session_state:
    st.session_state.library = {}
    st.session_state.messages = []
    st.session_state.document_understanding = DocumentUnderstanding()

tab1, tab2 = st.tabs(["Notes Generator", "Chat"])

with tab1:
# --- File Upload ---
    uploaded_file = st.file_uploader("Upload PDF/DOCX/PPTX", type=["pdf", "docx", "pptx"])

    if uploaded_file and uploaded_file.name not in st.session_state.library:
        with st.spinner("🔍 Analyzing document..."):
            try:
                if uploaded_file.type == "application/pdf":
                    pages = extract_pdf(uploaded_file)
                elif uploaded_file.type.endswith("document"):
                    pages = extract_docx(uploaded_file)
                elif uploaded_file.type.endswith("presentation"):
                    pages = extract_pptx(uploaded_file)
                else:
                    st.error("Unsupported file type")
                    st.stop()

                # Process document for understanding
                st.session_state.document_understanding.process_document(pages, compiled_model, tokenizer)

                # Generate notes (original functionality)
                sections = extract_headings_and_content(pages)
                notes_data = {}
                for heading, content in sections.items():
                    notes = generate_quality_notes(content)
                    if notes:
                        notes_data[heading] = notes
                st.session_state.library[uploaded_file.name] = notes_data

            except Exception as e:
                st.error(f"Error processing document: {e}")
                st.stop()

    # --- Notes Display ---
    if st.session_state.library:
        doc_names = list(st.session_state.library.keys())

        col3, col4 = st.columns(2, vertical_alignment="bottom", gap="small")

        with col3:
            selected_doc = st.selectbox("📚 Select Document", doc_names)
        search_term = st.text_input("🔍 Search in Notes")

        notes_data = st.session_state.library[selected_doc]
        all_notes = [note for bullets in notes_data.values() for note in bullets]

        with col4:
            st.download_button(
                "📥 Download as PDF",
                data=generate_pdf(all_notes),
                file_name=f"{selected_doc}_notes.pdf",
                mime="application/pdf", use_container_width=True
            )


        with st.expander("View All Notes", expanded=True):
            st.subheader("📝 Key Points")
            for head, bullets in notes_data.items():
                show = []
                for b in bullets:
                    if search_term.lower() in b.lower():
                        show.append(f"**{b}**")
                    elif not search_term:
                        show.append(b)
                if show:
                    st.markdown(f"### {head}")
                    for line in show:
                        st.markdown(line)

with tab2:
    if st.session_state.library:
        doc_names = list(st.session_state.library.keys())

    col1, col2 = st.columns([0.92, 0.08], vertical_alignment="bottom", gap="small")
    with col2:
    # Voice input
        if st.button("🎤"):
            query = listen()
            if query:
                st.session_state.last_question = query

    with col1:
    # Text input
        query_text = st.text_input(label="",placeholder="Ask Anything",
                             value=st.session_state.get('last_question', ''))

    if query_text:
        # Add to chat history
        st.session_state.messages.append({"role": "user", "content": query_text})

        # Get answer from document understanding system
        answer = st.session_state.document_understanding.answer_question(
            query_text, compiled_model, tokenizer
        )

        # Format answer
        formatted_answer = f"From the document:\n\n{answer}"
        st.session_state.messages.append({"role": "assistant", "content": formatted_answer})

    # Display chat history (most recent at top)
    with st.container(border=True):
        for message in reversed(st.session_state.messages[-5:]):
            with st.chat_message(message["role"]):
                st.markdown(message["content"])